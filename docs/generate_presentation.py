# -*- coding: utf-8 -*-
"""박유정 PetCare 발표 PPT 생성 스크립트 (상세 버전)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUTPUT = r"D:\petcare\docs\PetCare_발표_박유정.pptx"
OUTPUT_EN = r"D:\petcare\docs\PetCare_Presentation_ParkYujeong.pptx"

PRIMARY = RGBColor(0x2D, 0x6A, 0x4F)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY = RGBColor(0x88, 0x88, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x40, 0x9E, 0x6F)
CODE_BG = RGBColor(0xF4, 0xF6, 0xF8)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _header_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.07))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()


def _add_title_text(slide, title, top=Inches(0.25)):
    tx = slide.shapes.add_textbox(Inches(0.55), top, Inches(8.9), Inches(0.65))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.2), Inches(10), Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(8.4), Inches(1.5))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(8.4), Inches(3.0))
        tf2 = tx2.text_frame
        tf2.word_wrap = True
        for i, line in enumerate(subtitle.split("\n")):
            p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p2.text = line
            p2.font.size = Pt(15)
            p2.font.color.rgb = RGBColor(0xE0, 0xF0, 0xE8)
            p2.space_after = Pt(5)


def add_section_slide(prs, part_num, part_title, desc=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(2.5))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = part_num
    p.font.size = Pt(13)
    p.font.color.rgb = ACCENT
    p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = part_title
    p2.font.size = Pt(30)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    if desc:
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(14)
        p3.font.color.rgb = RGBColor(0xBB, 0xCC, 0xBB)
        p3.space_before = Pt(12)


def add_bullets(prs, title, bullets, font_main=12, font_sub=11, top=Inches(0.95)):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    _header_bar(slide)
    _add_title_text(slide, title)
    body = slide.shapes.add_textbox(Inches(0.55), top, Inches(8.9), Inches(6.2))
    tf = body.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = text
        para.level = min(level, 2)
        size = font_main if level == 0 else font_sub
        para.font.size = Pt(size - level)
        para.font.color.rgb = DARK if level == 0 else GRAY
        para.space_after = Pt(3)
        para.space_before = Pt(1)


def add_two_col(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    _header_bar(slide)
    _add_title_text(slide, title)

    for col_title, items, left in [(left_title, left_items, Inches(0.55)), (right_title, right_items, Inches(5.1))]:
        tx = slide.shapes.add_textbox(left, Inches(0.95), Inches(4.3), Inches(0.4))
        p = tx.text_frame.paragraphs[0]
        p.text = col_title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = ACCENT
        body = slide.shapes.add_textbox(left, Inches(1.35), Inches(4.3), Inches(5.8))
        tf = body.text_frame
        tf.word_wrap = True
        for i, text in enumerate(items):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = text
            para.font.size = Pt(11)
            para.font.color.rgb = DARK
            para.space_after = Pt(4)


def add_code_slide(prs, title, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    _header_bar(slide)
    _add_title_text(slide, title)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.95), Inches(9.0), Inches(6.0))
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(1.1), Inches(8.6), Inches(5.7))
    tf = tx.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        para.font.size = Pt(10)
        para.font.name = "Consolas"
        para.font.color.rgb = DARK
        para.space_after = Pt(2)


def add_table_slide(prs, title, headers, rows, font_size=9):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    _header_bar(slide)
    _add_title_text(slide, title)
    cols = len(headers)
    tbl_rows = len(rows) + 1
    row_h = min(0.32, 5.5 / tbl_rows)
    table = slide.shapes.add_table(tbl_rows, cols, Inches(0.45), Inches(0.95), Inches(9.1), Inches(row_h * tbl_rows)).table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(font_size + 1)
            para.font.bold = True
            para.font.color.rgb = WHITE
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(font_size)
                para.font.color.rgb = DARK


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ═══════════════════════════════════════
    # INTRO
    # ═══════════════════════════════════════
    add_title_slide(prs, "PetCare 반려동물 통합 플랫폼",
        "담당 모듈 개발 발표\n\n발표자: 박유정\n"
        "기술 스택: Spring Boot 3 · JSP/JSTL · MyBatis · Oracle\n"
        "외부 연동: Chart.js · Kakao Map API · 공공데이터 API\n"
        "작업 기간: 2026.07.06 ~ 2026.07.31 (약 4주)")

    add_bullets(prs, "발표 개요", [
        "① 프로젝트 소개 및 담당 범위",
        "② 가족찾기(Give) — 유기동물 API · 발견 신고 · 재능나눔",
        "③ 커뮤니티 — 게시글 · 댓글 · 좋아요 · 신고 · 수의사상담",
        "④ 관리자(Admin) — 대시보드 · 통계 · 회원 · 커뮤니티 · 리뷰",
        "⑤ 회원/인증 — 정지·탈퇴 정책 · Interceptor · 스케줄러",
        "⑥ 리뷰 연동 — 사업자 삭제요청 → 관리자 승인 · 알림 · 포인트",
        "⑦ 트러블슈팅 · 시연 · Q&A",
    ])

    add_table_slide(prs, "담당 모듈 전체 맵 (코드 주석 '박유정' 기준 150+ 파일)",
        ["대분류", "핵심 기능", "대표 URL", "작업 기간"],
        [
            ("가족찾기 Give", "APMS API·발견신고·카카오맵·재능나눔", "/give/*", "7/06~14, 7/29"),
            ("커뮤니티", "CRUD·댓글·좋아요·신고·LIFE상담", "/community/*", "7/08~10, 7/16"),
            ("관리자 Admin", "대시보드·통계·회원·커뮤니티·리뷰", "/admin/*", "7/15~31"),
            ("회원/인증", "정지·탈퇴·재가입차단·고객센터", "/member/*", "7/21~28"),
            ("사업자 Biz", "병원/숙소 리뷰·삭제요청·재능신청", "/biz/*", "7/14, 7/24~28"),
            ("마이페이지", "회원정보·숙소리뷰·포인트·알림", "/mypage/*", "7/28~29"),
        ], font_size=8)

    add_bullets(prs, "작업 타임라인 (주차별)", [
        "7월 1주 (7/06~07)  Give 모듈 착수 — 유기동물 공공 API 연동, 발견 신고 CRUD",
        "7월 2주 (7/08~14)  커뮤니티 전 기능 + 재능나눔(사용자·사업자·관리자 승인)",
        "7월 3주 (7/15~22)  관리자 커뮤니티·회원관리 + 정지/탈퇴 인증 체계 구축",
        "7월 4주 (7/24~31)  리뷰삭제요청·숙소리뷰·알림 + 대시보드·통계 실데이터 연동",
        "",
        "팀 공통 패턴: Controller → Service → Mapper(XML) → JSP",
        "주석 규칙: // 2026-07-XX 박유정 — 설명 (활동확인서·발표 추적용)",
    ])

    add_code_slide(prs, "공통 아키텍처 & 레이어 역할", [
        "[요청 흐름]",
        "Browser → @Controller → @Service → @Mapper 인터페이스",
        "                              ↓",
        "                    Mapper.xml (SQL) → Oracle DB",
        "                              ↓",
        "                    Model attribute → JSP 렌더링",
        "",
        "[레이어별 책임]",
        "Controller : URL 매핑, HttpSession, 로그인 체크, Model 전달만",
        "Service    : 비즈니스 로직, VO 조합, 빈 데이터 0 패딩, CSV 생성",
        "Mapper.xml : SELECT/INSERT/UPDATE/DELETE SQL만",
        "JSP        : JSTL c:forEach, Chart.js, 카카오맵 스크립트",
        "",
        "[외부 연동]",
        "공공데이터 API (유기동물) | KakaoMapService (지도·Geocode)",
        "Chart.js (대시보드·통계 차트) | @Cacheable (API 캐시)",
    ])

    # ═══════════════════════════════════════
    # PART 1: GIVE
    # ═══════════════════════════════════════
    add_section_slide(prs, "PART 1", "가족찾기 (Give) 모듈",
        "유기동물 공공 API · 시민 발견 신고 · 재능나눔 | 헤더 메뉴 '가족찾기'")

    add_bullets(prs, "Give 모듈 구조", [
        "모듈 진입: GET /give → redirect:/give/animal/list (GiveBaseController)",
        "",
        "① /give/animal  — 유기동물 조회 (정부 APMS 공공 API, DB 미사용)",
        "② /give/report  — 유기동물 발견 신고 (TB_POST + TB_FILE + 카카오맵)",
        "③ /give/talent  — 재능나눔 (TB_TALENT, 관리자 승인 후 노출)",
        "",
        "공통 탭 UI: give/index.jsp (animal / report / talent 탭 전환)",
        "7/16: 재능나눔 탭을 community/index.jsp에서 분리 → Give 모듈로 이동",
    ])

    add_two_col(prs, "유기동물 조회 — 구현 상세",
        "화면·UX", [
            "URL: /give/animal/list, /detail?desertionNo=",
            "JSP: give/animal/list.jsp, detail.jsp",
            "필터: 시도(sido), 축종(upkind), 보호상태(state)",
            "보호상태: 전체 / 보호중(protect) / 입양완료(adopt)",
            "페이징: 20건/페이지, 5개 단위 페이지 버튼",
            "데이터 범위: 최근 30일 (SEARCH_DAYS=30)",
            "첫 진입(search=false): API 호출 안 함 → 안내만",
            "[조회] 클릭(search=true): 그때 API 호출",
        ],
        "백엔드·API", [
            "GiveAnimalController → GiveAnimalServiceImpl",
            "API: abandonmentPublicService_v2 (data.go.kr)",
            "VO: AbandonmentVO (v2 필드 upKindCd·kindNm 파싱)",
            "결과: GiveAnimalListResult (animals, totalCount, totalPages)",
            "@Cacheable(value='animalList') — 동일 조건 재조회 캐시",
            "CacheConfig.java — @EnableCaching 활성화",
            "Controller에 있던 API 코드 → Service로 리팩터링",
            "JSON 파싱: Jackson ObjectMapper → List<AbandonmentVO>",
        ])

    add_code_slide(prs, "유기동물 API 호출 흐름 (GiveAnimalController)", [
        "@GetMapping('/list')",
        "public String animalList(sido, upkind, state, pageNo, search, model) {",
        "    if (!search) {",
        "        // 첫 진입: 빈 목록 + searched=false → JSP 안내문만 표시",
        "        model.addAttribute('searched', false);",
        "        return 'give/animal/list';",
        "    }",
        "    // [조회] 클릭 후: Service → 정부 API → 캐시 저장",
        "    GiveAnimalListResult result = giveAnimalService.getAnimalList(...);",
        "    model.addAttribute('animals', result.getAnimals());",
        "    model.addAttribute('totalCount', result.getTotalCount());",
        "    model.addAttribute('searched', true);",
        "    return 'give/animal/list';",
        "}",
        "",
        "상세: getAnimalDetail(desertionNo) → AbandonmentVO 1건 → detail.jsp",
    ])

    add_bullets(prs, "유기동물 발견 신고 — 기능 상세", [
        "목적: 시민이 발견한 유기동물 정보 등록 → 주인 찾기·보호소 연계",
        "URL: GET/POST /give/report/write | GET /give/report/list | /detail",
        "",
        "등록 흐름:",
        ("  1. GET /write → 로그인 체크 → write.jsp (비로그인 시 안내+로그인 링크)", 1),
        ("  2. POST /write → GiveReportServiceImpl.insertReport(vo, member, photos)", 1),
        ("  3. TB_POST INSERT (BOARD_TYPE='LOST') + 사진 저장 + TB_FILE INSERT", 1),
        ("  4. redirect → /give/report/detail?id={postId}", 1),
        "",
        "입력 항목: 동물종·품종·성별·크기·색상·특징태그·발견일시·발견위치·임시보호여부",
        "댓글: CommunityCommentService 재사용 (댓글·대댓글·수정·삭제)",
        "상태 변경: POST /give/report/status → 보호중/입양완료 등",
    ])

    add_two_col(prs, "발견 신고 — 카카오맵 연동 (7/29)",
        "저장 로직", [
            "GiveReportController — KakaoMapService DI",
            "write.jsp: 다음 주소검색 → Geocoder",
            "hidden: lostLat, lostLng → POST 전달",
            "GiveReportServiceImpl.insertReport():",
            "  · 폼 hidden 좌표 우선 사용",
            "  · 없으면 Kakao REST geocode(주소)",
            "  · 실패 시 서울시청 기본좌표 (37.5665, 126.978)",
            "TB_POST.LOST_LAT / LOST_LNG 저장",
            "common/kakaomap.jsp 공통 컴포넌트 재사용",
        ],
        "화면 표시", [
            "write.jsp: 지도 클릭·주소검색 시 마커 갱신",
            "detail.jsp: 저장된 LOST_LAT/LNG로 지도 표시",
            "Controller → model: kakaoJsApiKey, mapLat, mapLng",
            "bizName 속성 = 마커 위 라벨 (region명)",
            "병원 지도(BizHospital)와 동일 패턴 적용",
            "",
            "파일:",
            "GiveReportController.java",
            "GiveReportServiceImpl.java",
            "give/report/write.jsp, detail.jsp",
        ])

    add_table_slide(prs, "재능나눔 — 3자 연동 흐름",
        ["단계", "역할", "URL / 파일", "DB 상태"],
        [
            ("1 신청", "사업자(병원)", "POST /biz/hospital/talent", "STATUS_CD=PENDING"),
            ("2 검토", "관리자", "/admin/biz/talent", "승인→APPROVED / 반려→REJECTED"),
            ("3 노출", "일반 사용자", "/give/talent/list", "APPROVED만 조회"),
            ("배지", "관리자 사이드바", "AdminSidebarAdvice", "pendingTalentApproveCount"),
        ], font_size=8)

    add_bullets(prs, "재능나눔 — 구현 파일", [
        "GiveTalentController — 사용자 목록/상세 (APPROVED 필터)",
        "GiveTalentServiceImpl — applyTalent() 사업자 신청, 승인/반려 처리",
        "GiveTalentMapper.xml — INSERT, STATUS_CD 갱신, 사업자 이력 조회",
        "BizHospitalController — POST /biz/hospital/talent (talentType=HOSPITAL 고정)",
        "AdminBizController — POST /talent/approve, /talent/reject",
        "신청 전 검증: TB_BUSINESS STATUS_CD=APPROVED 여부 확인",
        "JSP: give/talent/list.jsp, detail.jsp | biz/hospital/talent.jsp | admin/biz/talent.jsp",
    ])

    # ═══════════════════════════════════════
    # PART 2: COMMUNITY
    # ═══════════════════════════════════════
    add_section_slide(prs, "PART 2", "커뮤니티 모듈",
        "게시글 · 댓글 · 좋아요 · 신고 · 수의사상담(LIFE) | 7/08~10")

    add_table_slide(prs, "커뮤니티 기능별 상세",
        ["기능", "구현", "테이블", "비고"],
        [
            ("게시글 목록", "CommunityPostServiceImpl.getPostList()", "TB_POST", "5건/페이지, thumbUrl"),
            ("게시글 등록", "insertPost() + 사진 최대 5장", "TB_POST, TB_FILE", "REF_TYPE='POST'"),
            ("게시글 상세", "getPostDetail() + VIEW_COUNT+1", "TB_POST", "photoUrls 목록"),
            ("댓글", "CommunityCommentServiceImpl", "TB_POST_COMMENT", "대댓글 parentId"),
            ("댓글 수정", "updateComment() — 작성자 검증", "TB_POST_COMMENT", "7/14 추가"),
            ("좋아요", "CommunityReactionServiceImpl", "TB_POST_REACTION", "토글 INSERT/DELETE"),
            ("신고", "CommunityReportServiceImpl", "TB_POST_REPORT", "게시글·댓글 대상"),
            ("LIFE상담", "TAGS=WAITING→ANSWERED", "TB_POST", "답변완료 박스 표시"),
        ], font_size=8)

    add_bullets(prs, "커뮤니티 — 수의사상담(LIFE) 특수 처리", [
        "게시판 타입: BOARD_TYPE = 'LIFE'",
        "등록 시: TAGS = 'WAITING' (답변 대기)",
        "수의사 답변 후: markLifeAnswered() → TAGS = 'ANSWERED'",
        "목록: LIFE 탭에서 글쓰기 버튼 숨김 (수의사만 답변)",
        "상세: ANSWERED 시 수의사 답변 박스 표시 (detail.jsp)",
        "목록 미리보기: 첫 일반댓글을 답변 미리보기로 표시 (STEP 4)",
        "",
        "7/16 UI 변경: community/index.jsp 탭 구성 변경",
        "  · 재능나눔 탭 제거 → give 모듈로 이동",
        "  · 전체 · 집사생활 · 무료나눔 · 수의사상담 4탭 유지",
    ])

    add_bullets(prs, "커뮤니티 — 핵심 파일 목록", [
        "Controller: CommunityPostController, CommunityCommentController, CommunityReactionController",
        "Service: CommunityPostServiceImpl, CommunityCommentServiceImpl, CommunityReactionServiceImpl",
        "Service: CommunityReportServiceImpl",
        "Mapper: CommunityPostMapper.xml, CommunityCommentMapper.xml, CommunityReactionMapper.xml",
        "VO: CommunityPostVO, CommunityCommentVO, CommunityReactionVO, CommunityReportVO",
        "JSP: community/index.jsp, list.jsp, detail.jsp, write.jsp",
        "신고 테이블명: TB_POST_REPORT (AdminCommunityVO 7/10 반영)",
    ])

    # ═══════════════════════════════════════
    # PART 3: ADMIN
    # ═══════════════════════════════════════
    add_section_slide(prs, "PART 3", "관리자 (Admin) 모듈",
        "대시보드 · 통계 · 회원 · 커뮤니티 · 리뷰 · 재능승인 | 7/15~31")

    add_bullets(prs, "ADMIN-01 대시보드 (/admin) — 개요", [
        "Controller: AdminMainController.dashboard()",
        "Service: AdminMainServiceImpl.getDashboardSummary()",
        "VO: AdminMainVO (통계·차트·목록 데이터 한 번에 담음)",
        "JSP: admin/dashboard.jsp + Chart.js",
        "",
        "기존 더미 데이터(하드코딩 3행) 전부 제거 → Mapper SQL 실시간 집계로 교체",
        "Phase 1 (7/29): 승인 대기 사업자 5건 (TB_BUSINESS PENDING)",
        "Phase 2 (7/30): 상단 통계 카드 4종",
        "Phase 3 (7/30): 최근 주문·회원 도넛·매출 차트",
    ])

    add_table_slide(prs, "대시보드 — 상단 통계 카드 4종",
        ["카드", "SQL 메서드", "테이블·조건", "JSP 표시"],
        [
            ("오늘 가입", "countTodayJoin()", "TB_MEMBER JOIN_DATE=오늘", "summary.todayJoinCount"),
            ("오늘 주문", "countTodayOrder()", "TB_ORDER ORDER_DATE=오늘", "summary.todayOrderCount"),
            ("오늘 매출", "sumTodaySales()", "PAY_AMOUNT, 취소제외", "원→백만원 ÷1000000"),
            ("미처리 예약", "countPendingResv()", "PENDING+CONFIRMED", "summary.pendingResvCount"),
        ], font_size=8)

    add_two_col(prs, "대시보드 — 차트 구현 상세",
        "매출 차트 (Phase 3-C)", [
            "주간: 최근 7일 일별 매출 (요일 라벨)",
            "월간: 이번 달 1일~말일 일별 매출",
            "switchSalesChart() JS 함수로 전환",
            "Service: DB 결과 → Map(라벨→금액)",
            "주문 없는 날 = 0원 패딩 (for문으로 7일/말일 채움)",
            "JSP: 원 → 만원 (÷10000) Chart.js line chart",
            "SQL: ORDER_STATUS <> 'CANCEL' 조건",
            "ORDER_DATE: DATE 타입, TRUNC() 비교",
        ],
        "회원 도넛 (Phase 3-B)", [
            "3분류: 일반 / 사업자 / 탈퇴",
            "일반: 탈퇴·승인사업자 제외 COUNT",
            "사업자: TB_BUSINESS STATUS=APPROVED",
            "탈퇴: STATUS_CD='WITHDRAWN'",
            "JSP: 합계로 % 계산 → Chart.js doughnut",
            "범례 3줄: 건수 + 비율% 표시",
            "",
            "최근 주문 5건 (Phase 3-A):",
            "TB_ORDER + TB_MEMBER JOIN, ORDER_DATE DESC",
        ])

    add_bullets(prs, "ADMIN-04 통계 (/admin/stats) — Phase별 상세", [
        "Controller: AdminMainController.stats() → getStatsSummary()",
        "VO: AdminStatsVO (요약·차트·전월대비·업종별 데이터)",
        "JSP: admin/stats/index.jsp + Chart.js 3개",
        "",
        "Phase 1: 이번 달 매출(취소제외)·가입·예약·주문 4카드",
        "Phase 2: 월별 매출 line chart — 최근 6개월, Y축 백만원",
        "Phase 3: 월별 신규 가입 bar chart — 최근 6개월",
        "Phase 4: 업종별 예약/주문 — 병원·숙소·쇼핑 UNION ALL 3행",
        "Phase 5-A: 전월 대비 % — calcChangeRate() → JSP c:choose ▲/▼",
        "Phase 5-C: GET /admin/stats/export — CSV 4섹션 + BOM(한글)",
        "Phase 5-B: 기간 필터 — 보류 (UI만 존재, 로직 미구현)",
    ])

    add_code_slide(prs, "통계 — 6개월 패딩 로직 (Service)", [
        "// Phase 2: 월별 매출 — DB에 없는 달도 0원으로 채움",
        "LocalDate start = LocalDate.now().minusMonths(5).withDayOfMonth(1);",
        "for (int i = 0; i < 6; i++) {",
        "    String key = start.format(DateTimeFormatter.ofPattern('yyyy-MM'));",
        "    String label = start.getMonthValue() + '월';  // '3월', '4월' ...",
        "    long amount = salesMap.getOrDefault(key, 0L);",
        "    // AdminMainSalesDayVO(dayLabel, salesAmount) 리스트에 add",
        "    start = start.plusMonths(1);",
        "}",
        "",
        "// Phase 5-A: 전월 대비 증감률",
        "calcChangeRate(current, previous):",
        "  previous==0 → current>0 이면 100%, 아니면 0%",
        "  else → (current-previous)/previous*100, 소수1자리",
    ])

    add_table_slide(prs, "통계 Phase 4 — 업종별 예약/주문 SQL",
        ["라벨", "데이터 소스", "집계 조건", "비고"],
        [
            ("병원", "TB_RESERVATION", "RESV_TYPE='HOSPITAL', REG_DATE 이번달", "예약 건수"),
            ("숙소", "TB_RESERVATION", "RESV_TYPE='STAY', REG_DATE 이번달", "예약 건수"),
            ("쇼핑", "TB_ORDER", "ORDER_DATE 이번달", "취소 포함"),
        ], font_size=9)

    add_bullets(prs, "통계 CSV Export (Phase 5-C)", [
        "URL: GET /admin/stats/export → AdminMainController.exportStatsCsv()",
        "Service: exportStatsCsv() — getStatsSummary() 재사용 (중복 SQL 없음)",
        "BOM: writer.write('\\uFEFF') — Excel 한글 깨짐 방지",
        "섹션 1: 요약 카드 4항목 + 전월대비 %",
        "섹션 2: 월별 매출 추이 (최근 6개월, 원 단위)",
        "섹션 3: 월별 신규 가입자 (최근 6개월)",
        "섹션 4: 업종별 예약/주문 (병원·숙소·쇼핑)",
        "CSV 이스케이프: 쉼표·따옴표 포함 시 \" 감싸기 (회원 export 동일 패턴)",
        "파일명: stats_YYYYMMDD.csv",
    ])

    add_bullets(prs, "관리자 회원 관리 — STEP별 상세", [
        "URL: /admin/member/list, /detail, POST suspend/restore/withdraw ...",
        "Controller: AdminMemberController | Service: AdminMemberServiceImpl",
        "Mapper: AdminMemberMapper.xml | JSP: list.jsp, detail.jsp",
        "",
        "STEP 7 (7/20): 정지·복구·강제탈퇴 — STATUS_CD 변경",
        ("  · 정지 유형: DAY3 / DAY7 / PERMANENT → SUSPEND_END_DATE 저장", 1),
        ("  · 복구: STATUS_CD=NORMAL, SUSPEND_END_DATE=NULL", 1),
        ("  · 강제탈퇴: STATUS_CD=WITHDRAWN (TB_MEMBER만, WITHDRAW 테이블 미INSERT)", 1),
        "STEP 8: 상세 화면 활동현황 — 주문·예약·게시글·신고·쿠폰·찜·펫 수",
        "STEP 9: 등급 변경 (GRADE_CD) | STEP 10: 포인트 적립·차감 + TB_POINT 이력",
        "STEP 11: 최근 주문 5건 | STEP 12: 체크박스 일괄 정지·탈퇴·복구",
        "CSV export: GET /admin/member/export — 검색조건 반영, BOM 포함",
        "스케줄러: MemberSuspendScheduler — 매일 0시 releaseExpiredSuspensions()",
    ])

    add_bullets(prs, "관리자 커뮤니티 관리 (7/15)", [
        "URL: /admin/community/list, /detail",
        "Service: AdminCommunityServiceImpl",
        "목록: 키워드·게시판타입·상태 필터 + 신고 건수 JOIN",
        "상세: 게시글 본문 + 사진 + 댓글 목록",
        "",
        "상태 처리 (소프트 삭제):",
        ("  · 숨김: STATUS_CD = 'HIDDEN'", 1),
        ("  · 삭제: STATUS_CD = 'DELETED'", 1),
        ("  · 복구: STATUS_CD = 'ACTIVE'", 1),
        "상태별 버튼 분기: ACTIVE→숨김/삭제, HIDDEN/DELETED→복구",
        "JSP: admin/community/list.jsp, detail.jsp",
    ])

    add_bullets(prs, "관리자 리뷰 삭제요청 (7/24~28)", [
        "URL: /admin/review/list — GET 검색(keyword, statusCd)",
        "대상: HOSPITAL + STAY + PRODUCT(쇼핑) 통합 — AdminReviewMapper.xml",
        "승인: AdminReviewApproveTxService — DELETE + 상태변경 (짧은 TX)",
        "       → AdminReviewSideEffectService — 평점 갱신·알림 (별도 TX)",
        "반려: REJECT_REASON 저장 + sendReviewDeleteRejectNotification()",
        "사이드바: AdminSidebarAdvice.pendingReviewDeleteCount",
        "JSP: admin/review/list.jsp — 승인 fetch(AJAX), 20초 타임아웃",
        "리뷰 타입별 테이블 분기: reviewType STAY/HOSPITAL/PRODUCT",
    ])

    # ═══════════════════════════════════════
    # PART 4: MEMBER AUTH
    # ═══════════════════════════════════════
    add_section_slide(prs, "PART 4", "회원/인증 & 정지 처리",
        "정지·탈퇴 정책 · Interceptor · 스케줄러 · 재가입 차단 | 7/21~28")

    add_table_slide(prs, "정지(SUSPENDED) vs 탈퇴(WITHDRAWN) 정책 비교",
        ["항목", "정지 SUSPENDED", "탈퇴 WITHDRAWN"],
        [
            ("로그인", "허용 (이메일/카카오 모두)", "차단 — MemberLoginBlockedException"),
            ("로그인 후", "/member/cs 고객센터로 redirect", "login.jsp 에러 메시지"),
            ("사업자 권한", "부여 안 함 (status 체크)", "—"),
            ("URL 접근", "Interceptor → /member/cs만", "—"),
            ("재가입", "가능", "이메일 영구 차단 (countWithdrawnMemberByEmail)"),
            ("DB 변경", "STATUS_CD + SUSPEND_END_DATE", "STATUS_CD만 (관리자 강제 시)"),
            ("자동 복구", "스케줄러 + 로그인 시 releaseExpired", "7일 후 purge (자발 탈퇴만)"),
        ], font_size=8)

    add_bullets(prs, "정지 회원 처리 — 구현 파일 & 흐름", [
        "① 관리자 정지: AdminMemberServiceImpl.suspendMember(memberNo, suspendType)",
        "   → DAY3: +3일, DAY7: +7일, PERMANENT: SUSPEND_END_DATE=NULL",
        "",
        "② 로그인: MemberAuthServiceImpl — WITHDRAWN만 차단, SUSPENDED는 허용",
        "   → sessionMember.setStatus('SUSPENDED') 세션에 저장",
        "   → MemberAuthController: SUSPENDED면 /member/cs redirect",
        "",
        "③ URL 제한: SuspendedMemberInterceptor (WebConfig 등록)",
        "   → 허용: /member/cs/**, /member/logout, /resources/, /upload/, cart/noti Ajax",
        "   → 그 외: /member/cs?restricted=1 redirect",
        "",
        "④ UI: header.jsp 정지 안내 배너 | member/cs.jsp 안내 화면",
        "⑤ 자동 복구: MemberSuspendScheduler (@Scheduled cron='0 0 0 * * *')",
        "   PetcareApplication @EnableScheduling 활성화",
    ])

    add_bullets(prs, "회원가입·인증 추가 작업", [
        "탈퇴 이메일 재가입 차단: countWithdrawnMemberByEmail() — GET/POST 모두 검증",
        "이메일 중복 확인: WITHDRAWN 회원 제외 (MemberAuthMapper.xml)",
        "생년월일·성별 저장: join.jsp → TB_MEMBER BIRTH_DATE, GENDER (7/28)",
        "MemberRegisterVO, MemberAuthServiceImpl — birthDate 필수 검증",
        "login.jsp: 탈퇴 회원 전용 에러 메시지 (정지는 로그인 후 cs)",
        "MemberLoginBlockedException — errorCode → login.jsp error 파라미터 매핑",
    ])

    # ═══════════════════════════════════════
    # PART 5: REVIEW & BIZ
    # ═══════════════════════════════════════
    add_section_slide(prs, "PART 5", "리뷰 & 사업자 연동",
        "삭제요청 · 숙소리뷰 · 포인트 · 알림 | 7/24~29")

    add_bullets(prs, "리뷰 삭제요청 — End-to-End 흐름", [
        "[사업자 요청]",
        "병원: BizHospitalController POST /reviews/delete-request",
        "숙소: BizStayController POST /reviews/delete-request",
        "→ TB_REVIEW_DELETE_REQUEST INSERT (STATUS_CD=PENDING)",
        "→ 동일 리뷰 PENDING 중복 요청 방지 (countPendingReviewDeleteRequest)",
        "",
        "[관리자 처리]",
        "AdminReviewController — /admin/review/list",
        "승인: TB_REVIEW DELETE + 요청 STATUS=APPROVED + 평점 갱신",
        "반려: REJECT_REASON 저장 + STATUS=REJECTED",
        "",
        "[알림]",
        "승인: sendReviewDeleteApproveNotification() → 사업자",
        "반려: sendReviewDeleteRejectNotification() → 사업자",
    ])

    add_two_col(prs, "사업자 숙소 리뷰관리 (7/27~28)",
        "기능", [
            "URL: /biz/stay/reviews",
            "탭 1: 리뷰 목록 — 답글 작성/수정",
            "탭 2: 삭제요청 이력 — 상태·처리일",
            "POST /reviews/reply — TB_REVIEW BIZ_REPLY",
            "POST /reviews/delete-request — 사유 필수",
            "PENDING 건: 답글·삭제요청 버튼 숨김",
            "Controller → JSON 변환 → reviews.jsp JS 렌더링",
            "사이드바 배지: PENDING+CONFIRMED 예약 건수",
        ],
        "파일", [
            "BizStayController.java",
            "BizStayServiceImpl.java",
            "BizStayMapper.xml",
            "biz/stay/reviews.jsp",
            "",
            "병원 동일 패턴:",
            "BizHospitalController / Service / Mapper",
            "biz/hospital/reviews.jsp",
            "진료기록 MEMO 파싱:",
            "MedicalRecordMemoParser — 유형·신체계측 분리",
        ])

    add_bullets(prs, "숙소 리뷰 · 포인트 · 알림 (7/28~29)", [
        "사용자 리뷰: MypageReserveServiceImpl.addStayReview()",
        "  → TB_REVIEW INSERT → updateStayRatingSummary() (AVG_RATING, REVIEW_CNT)",
        "  → sendStayReviewToBizNotification() 사업자 알림",
        "  → MypageReservePointService — 결제금액 3% 포인트 (별도 TX, 락 방지)",
        "  → StayReviewRegisterResult 반환 → Controller 세션 포인트 갱신",
        "",
        "사업자 답글: sendStayReviewReplyNotification() → 리뷰 작성 회원",
        "예약 결제: StayServiceImpl — sendStayReserveToBizNotification()",
        "숙소 목록: stay/list.jsp — AVG_RATING, REVIEW_CNT 표시",
        "mypage/reserve-detail.jsp — 포인트 적립 완료 안내 메시지",
    ])

    # ═══════════════════════════════════════
    # PART 6: MISC
    # ═══════════════════════════════════════
    add_section_slide(prs, "PART 6", "부가 기능 & 공통 연동")

    add_bullets(prs, "마이페이지 · 펫 프로필 · 기타", [
        "회원정보 수정: MypageAccountController /edit — DB 최신 프로필 조회",
        "MypageAccountVO: birthDate, gender 필드 추가 (7/28)",
        "펫 프로필: PetProfileVO — furColor, neuterYn, traits, memo (TB_PET)",
        "pets.jsp — 반려동물 등록/수정 화면 필드 연동",
        "AdminSidebarAdvice — 재능승인·리뷰삭제·사업자승인 배지 통합",
        "footer.jsp — 헤더 뱃지 Ajax 갱신 (.header-utils)",
        "WebConfig — SuspendedMemberInterceptor 등록 경로 설정",
    ])

    add_table_slide(prs, "주요 테이블 & 담당 기능 매핑",
        ["테이블", "주요 컬럼", "담당 기능"],
        [
            ("TB_MEMBER", "STATUS_CD, SUSPEND_END_DATE, GRADE_CD", "회원관리·인증"),
            ("TB_ORDER", "ORDER_DATE, PAY_AMOUNT, ORDER_STATUS", "대시보드·통계 매출"),
            ("TB_RESERVATION", "RESV_TYPE, REG_DATE, STATUS_CD", "통계·예약배지"),
            ("TB_POST", "BOARD_TYPE, STATUS_CD, LOST_LAT/LNG", "커뮤니티·신고"),
            ("TB_REVIEW", "TARGET_ID, RATING, BIZ_REPLY", "리뷰·평점"),
            ("TB_REVIEW_DELETE_REQUEST", "REVIEW_TYPE, STATUS_CD", "삭제요청"),
            ("TB_TALENT", "TALENT_TYPE, STATUS_CD", "재능나눔"),
            ("TB_POINT", "POINT_TYPE, AMOUNT", "포인트 이력"),
        ], font_size=8)

    add_bullets(prs, "트러블슈팅 — 실제 이슈 & 해결", [
        "ORDER_DATE VARCHAR2 → DATE DDL 변경",
        ("  · 증상: ORA-00932 inconsistent datatypes", 1),
        ("  · 해결: TRUNC(ORDER_DATE)=TRUNC(SYSDATE) 비교로 통일", 1),
        "",
        "공공 API 응답 지연 (3~5초)",
        ("  · 해결: search=false 첫 진입 미호출 + @Cacheable", 1),
        "",
        "대시보드/통계 차트 빈 구간",
        ("  · 해결: Service에서 Map + for문 0 패딩 (7일/6개월)", 1),
        "",
        "리뷰 등록 + 포인트 적립 동시 TX → DB 락 대기",
        ("  · 해결: MypageReservePointService, RatingService 별도 TX", 1),
        "",
        "관리자 강제탈퇴 vs 자발탈퇴 TB_MEMBER_WITHDRAW 불일치",
        ("  · 현황: 관리자 강제탈퇴는 TB_MEMBER만 변경 → purge 대상 아님", 1),
        ("  · 제안: WITHDRAW_TYPE 컬럼 (SELF/ADMIN) 팀 논의 필요", 1),
    ])

    add_bullets(prs, "시연 시나리오 (권장 5분)", [
        "① 가족찾기: /give/animal/list → 서울 선택 → [조회] → 상세 → 보호소 정보",
        "② 발견신고: /give/report/write → 주소검색·지도 클릭 → 사진첨부 → 등록",
        "③ 대시보드: /admin → 카드 4종 확인 → 주간/월간 차트 전환 → 도넛 차트",
        "④ 통계: /admin/stats → 3개 차트 + 전월대비 % → [Excel보내기] CSV 확인",
        "⑤ 회원관리: /admin/member/list → 검색 → 상세 → 정지(7일) → 복구",
        "⑥ 리뷰: /admin/review/list → PENDING 건 → 승인(AJAX) → 사업자 알림",
        "",
        "스크린샷 권장: 위 6화면 캡처 → 슬라이드에 삽입",
    ])

    add_bullets(prs, "성과 요약", [
        "가족찾기(Give) 모듈 전체 구현 — 공공 API + 시민 제보 + 재능나눔 3자 연동",
        "커뮤니티 핵심 기능 (게시·댓글·좋아요·신고·LIFE상담) 완성",
        "관리자 대시보드·통계 — 더미 제거, 실데이터 SQL 집계 + Chart.js 시각화",
        "회원 관리 12 STEP + 정지/탈퇴 인증 체계 + 스케줄러 자동 복구",
        "리뷰 삭제요청 병원·숙소·쇼핑 통합 + 알림·포인트·평점 연동",
        "외부 API 3종 연동: APMS(유기동물), Kakao Map, Chart.js",
    ])

    add_title_slide(prs, "감사합니다",
        "Q & A\n\n박유정 | PetCare 프로젝트\n"
        "Git 브랜치: parkyujeong\n"
        "코드 추적: 프로젝트 검색 '박유정'")

    prs.save(OUTPUT)
    prs.save(OUTPUT_EN)
    print(f"생성 완료: {OUTPUT}")
    print(f"영문 복사: {OUTPUT_EN}")
    print(f"총 슬라이드: {len(prs.slides)}장")


if __name__ == "__main__":
    build()
